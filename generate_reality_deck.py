#!/usr/bin/env python
"""
Bus Sathi — Trace Intelligence briefing deck.

A standalone presentation explaining how REAL driver GPS from the Bus Sathi
app was turned into a ground-truth layer for the Kashmir route plan: the
pipeline, what it verified, what it measured, how it refined the plan
(v3.4.5), and — honestly — what it cannot do.

Output: Bus_Sathi_Trace_Intelligence_Briefing.pptx (repo root)
Run:  & "D:\\plotting\\ana\\python.exe" generate_reality_deck.py
"""
import sys
sys.stdout.reconfigure(encoding="utf-8", errors="replace")
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── palette: deep teal transit + slate dark, amber accent ─────────
INK    = RGBColor(0x10, 0x2A, 0x2E)   # near-black teal ink
TEAL   = RGBColor(0x0F, 0x6E, 0x56)   # primary (matches repo plots)
TEALDK = RGBColor(0x0A, 0x3D, 0x33)   # dark teal for dark slides
MIST   = RGBColor(0xE9, 0xF2, 0xEE)   # light teal tint (cards)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
AMBER  = RGBColor(0xD9, 0x77, 0x06)   # accent / caution
RED    = RGBColor(0xB0, 0x43, 0x2F)
GREY   = RGBColor(0x5B, 0x6B, 0x6B)
CARD_W = 2.86

W, H = Inches(13.333), Inches(7.5)
prs = Presentation(); prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]
F_HDR, F_BODY = "Cambria", "Calibri"


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    return s


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, color=INK, bold=False, font=F_BODY, first=None,
         align=PP_ALIGN.LEFT, space=4, italic=False):
    p = tf.paragraphs[0] if (first if first is not None else not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return p


def card(s, x, y, w, h, fill=MIST, line=None):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.adjustments[0] = 0.06
    c.fill.solid(); c.fill.fore_color.rgb = fill
    if line: c.line.color.rgb = line; c.line.width = Pt(1)
    else: c.line.fill.background()
    c.shadow.inherit = False
    return c


def stat(s, x, y, w, big, label, sub, big_color=TEAL, on_dark=False):
    card(s, x, y, w, 1.62, fill=(TEALDK if on_dark else MIST))
    tf = box(s, x + 0.18, y + 0.14, w - 0.36, 1.36)
    para(tf, big, 30, big_color, bold=True, font=F_HDR)
    para(tf, label, 12.5, WHITE if on_dark else INK, bold=True)
    para(tf, sub, 10, RGBColor(0xB9, 0xD6, 0xCB) if on_dark else GREY, space=0)


def title_bar(s, kicker, title, sub=None):
    tf = box(s, 0.55, 0.32, 12.2, 1.3)
    para(tf, kicker.upper(), 12, TEAL, bold=True)
    para(tf, title, 31, INK, bold=True, font=F_HDR, space=2)
    if sub: para(tf, sub, 13, GREY, space=0)


def pic(s, path, x, y, w=None, h=None):
    kw = {}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    return s.shapes.add_picture(path, Inches(x), Inches(y), **kw)


def caption(s, x, y, w, text, align=PP_ALIGN.LEFT):
    tf = box(s, x, y, w, 0.5)
    para(tf, text, 11.5, GREY, italic=True, align=align, space=0)


# ════════ 1. TITLE (dark) ═══════════════════════════════════════════
s = slide(TEALDK)
tf = box(s, 0.9, 1.6, 11.4, 3.6)
para(tf, "BUS SATHI  ·  TRACE INTELLIGENCE", 15, RGBColor(0x7F, 0xC8, 0xAF), bold=True)
para(tf, "Ground-Truthing the Kashmir Plan\nwith Real Driver GPS", 44, WHITE, bold=True, font=F_HDR, space=10)
para(tf, "What 3.8 million real GPS points from working bus drivers verified, measured,",
     16, RGBColor(0xD5, 0xE8, 0xDF), space=0)
para(tf, "and corrected in the route-rationalisation plan — and what they honestly cannot tell us.",
     16, RGBColor(0xD5, 0xE8, 0xDF), space=0)
tf2 = box(s, 0.9, 6.35, 11.4, 0.6)
para(tf2, "Kashmir Valley Route Rationalisation  ·  companion to plan v3.4.5  ·  July 2026", 12, RGBColor(0x9F, 0xBF, 0xB3))

# ════════ 2. THE DATA ═══════════════════════════════════════════════
s = slide()
title_bar(s, "the raw material", "Every 5 seconds, a real bus told us where it was",
          "The Bus Sathi Android app (published on Google Play) logs a driver's position every ~5 s into Firebase.")
xs = [0.55, 3.79, 7.03, 10.27]
stat(s, xs[0], 1.85, CARD_W, "180", "volunteer drivers", "self-registered via the public app")
stat(s, xs[1], 1.85, CARD_W, "3.8M", "GPS points", "~5-second sampling, Feb–Jun 2026")
stat(s, xs[2], 1.85, CARD_W, "1,213", "recorded sessions", "one session = hours of a working day")
stat(s, xs[3], 1.85, CARD_W, "131", "observed days", "weekday + weekend service")
tf = box(s, 0.55, 3.95, 12.2, 1.1)
para(tf, "The catch: a recorded “trip” is NOT a bus trip.", 17, INK, bold=True, font=F_HDR)
para(tf, "Drivers leave the app running in the background — one session bundles 3–10 real service runs plus hours "
         "of parking, meals and wandering. Used naively, this data gives absurd timings. So nothing is used raw.",
     13.5, GREY, space=0)
c = card(s, 0.55, 5.25, 12.2, 1.55, fill=MIST)
tf = box(s, 0.85, 5.45, 11.6, 1.2)
para(tf, "Privacy by construction", 13.5, TEAL, bold=True)
para(tf, "Driver identities are SHA-256 hashed in every derived file; raw traces and credentials never leave the "
         "private data store; only aggregate, anonymised layers are published.", 12.5, INK, space=0)

# ════════ 3. SESSIONISING ═══════════════════════════════════════════
s = slide()
title_bar(s, "step 1 — cleaning", "From messy sessions to honest service runs",
          "A dwell/gap segmenter cuts each session at terminals (≥12 min stationary) and app gaps (≥8 min), then trims idle heads and tails.")
stat(s, xs[0], 1.85, CARD_W, "2,526", "real service runs", "extracted from 1,200 sessions")
stat(s, xs[1], 1.85, CARD_W, "2,012 h", "idle time removed", "parking / breaks / background noise")
stat(s, xs[2], 1.85, CARD_W, "63 min", "median real run", "vs 209 min raw session wall-time")
stat(s, xs[3], 1.85, CARD_W, "13.3", "km/h median run speed", "realistic city-bus pace → sanity check passed")
tf = box(s, 0.55, 3.9, 12.2, 0.55)
para(tf, "Then every run is snapped to the road network (OSRM map-matching) and quality-gated:", 14, INK, bold=True)
rows = [
    ("2,519 / 2,526 runs matched to roads", "raw-vs-matched length agreement used as the gate — median 0.99"),
    ("2,426 runs pass as “clean”", "these are the only runs any downstream claim is built on"),
    ("38,431 stop-dwell events kept", "each 40 s – 12 min stop becomes evidence for the stop inventory"),
]
y = 4.55
for h, d in rows:
    card(s, 0.55, y, 12.2, 0.78, fill=MIST)
    tf = box(s, 0.85, y + 0.13, 11.6, 0.55)
    p = para(tf, h + "   ", 13.5, TEAL, bold=True)
    r = p.add_run(); r.text = d; r.font.size = Pt(12.5); r.font.color.rgb = GREY; r.font.name = F_BODY
    y += 0.92

# ════════ 3b. VISUAL — SESSIONISING ═════════════════════════════════
s = slide()
title_bar(s, "what the cleaning looks like", "One recorded session, ten real bus runs inside it")
pic(s, "data/session_split.png", 1.15, 1.55, w=11.0)
caption(s, 1.15, 7.02, 11.0,
        "Left: one raw Firestore session (7,906 points — a whole working day). Right: the same session decomposed "
        "into 10 idle-trimmed service runs, each with its own realistic distance and duration.", align=PP_ALIGN.CENTER)

# ════════ 3c. VISUAL — MAP-MATCHING ═════════════════════════════════
s = slide()
title_bar(s, "quality gate", "Raw GPS (red) vs road-matched (green)")
pic(s, "data/sample_before_after.png", 2.75, 1.42, w=7.85)
caption(s, 1.15, 6.72, 11.0,
        "On genuine runs the matched line hugs the raw trace. The top-right panel is parked-vehicle GPS scribble — "
        "exactly what the agreement gate rejects before any analysis.", align=PP_ALIGN.CENTER)

# ════════ 4. CORRIDORS + VERDICTS ═══════════════════════════════════
s = slide()
title_bar(s, "step 2 — verification", "18 observed corridors, each judged one by one",
          "Endpoint clustering (≥3 distinct drivers per terminal) grouped the clean runs into 25 corridors; the 18 with real support (≥5 runs, ≥2 drivers) each got an evidence file and an individual AI-analyst verdict.")
tf = box(s, 0.55, 2.0, 6.0, 0.4)
para(tf, "Verdicts vs the plan's routes", 15, INK, bold=True, font=F_HDR)
vd = [("7", "MATCH the plan", "same corridor, same terminals — e.g. Soura–Lal Chowk = FDR-050 (211 runs)", TEAL),
      ("8", "geometry diverges", "buses observably run a different alignment than the routed line → reconciliation list", AMBER),
      ("2", "outside Kashmir Division", "NH-44 / Jammu traffic — correctly excluded by the district clip", GREY),
      ("1", "data artifact", "near-stationary terminal idling; excluded and the filter fixed", GREY)]
y = 2.5
for n, h, d, col in vd:
    card(s, 0.55, y, 7.35, 1.0, fill=MIST)
    tf = box(s, 0.8, y + 0.12, 0.9, 0.8)
    para(tf, n, 27, col, bold=True, font=F_HDR)
    tf = box(s, 1.75, y + 0.14, 6.0, 0.8)
    para(tf, h, 13.5, INK, bold=True, space=1)
    para(tf, d, 11.5, GREY, space=0)
    y += 1.14
c = card(s, 8.2, 2.5, 4.55, 4.4, fill=TEALDK)
tf = box(s, 8.5, 2.75, 3.95, 3.9)
para(tf, "0", 48, WHITE, bold=True, font=F_HDR)
para(tf, "informal / unpermitted routes", 15, WHITE, bold=True, space=8)
para(tf, "An early “under-permitted cluster” finding was RETRACTED after checking the raw 614-permit register — "
         "those areas are heavily permitted. The honest signal is geometry divergence, not illegal operation.",
     12, RGBColor(0xC9, 0xE2, 0xD8), space=8)
para(tf, "The audit trail (verdict + correction per corridor) ships with the repo.", 11, RGBColor(0x9F, 0xBF, 0xB3), italic=True, space=0)

# ════════ 4b. VISUAL — OVERLAY EVIDENCE ═════════════════════════════
s = slide()
title_bar(s, "the overlay evidence", "Observed corridor (green) laid over the plan's candidate routes")
pic(s, "analyst/evidence/C1.png", 0.75, 1.6, w=5.35)
pic(s, "analyst/evidence/C2.png", 7.25, 1.6, w=5.35)
tf = box(s, 0.75, 6.98, 5.35, 0.5)
para(tf, "C1 — a MATCH: 211 observed runs ride exactly on permit FDR-050 (Soura–Lal Chowk).",
     11.5, TEAL, bold=True, align=PP_ALIGN.CENTER, space=0)
tf = box(s, 7.25, 6.98, 5.35, 0.5)
para(tf, "C2 — a DIVERGENCE: 150 real runs (green) vs the routed lines — the alignment to reconcile.",
     11.5, AMBER, bold=True, align=PP_ALIGN.CENTER, space=0)

# ════════ 5. MEASURED SPEEDS ════════════════════════════════════════
s = slide()
title_bar(s, "step 3 — measurement", "The first measured bus speeds for Srinagar",
          "Not modelled, not assumed — computed from consecutive GPS fixes of buses in revenue service.")
stat(s, 0.55, 1.85, 3.45, "21 km/h", "moving speed", "median while the bus is in motion")
stat(s, 4.2, 1.85, 3.45, "12.5 km/h", "effective speed", "including every service stop")
stat(s, 0.55, 3.62, 3.45, "37%", "of run time at stops", "boarding + waiting dwell share")
stat(s, 4.2, 3.62, 3.45, "17 vs 24", "core vs periphery km/h", "10,600-cell congestion heatmap")
c = card(s, 0.55, 5.45, 7.1, 1.55, fill=MIST)
tf = box(s, 0.8, 5.62, 6.6, 1.25)
para(tf, "Why it matters: the engine's drive times come from a CAR model (55–137 km/h “free flow” here). "
         "This layer replaces that assumption with measurement — per 100 m cell, per corridor.",
     12, INK, space=0)
pic(s, "data/speed_layer.png", 8.05, 1.55, h=5.45)
caption(s, 8.05, 7.04, 4.9, "Measured speed per cell — red = congested crawl.", align=PP_ALIGN.CENTER)

# ════════ 6. OPERATIONS ═════════════════════════════════════════════
s = slide()
title_bar(s, "step 4 — operations", "What a driver's day actually looks like",
          "855 measured driver-days — duty spans, utilisation, terminal turnarounds and the shape of the service day.")
stat(s, xs[0], 1.85, CARD_W, "7.9 h", "duty-day span", "4.8 h of it in service = 76% utilisation")
stat(s, xs[1], 1.85, CARD_W, "24 min", "median terminal turn", "door-to-door incl. layover (n=1,178)")
stat(s, xs[2], 1.85, CARD_W, "10:20–18:15", "the typical service day", "ramps from 8:00, peaks 14:00–17:00")
stat(s, xs[3], 1.85, CARD_W, "~19:00", "service collapses", "the observed fleet stops by evening")
c = card(s, 0.55, 3.95, 12.2, 2.15, fill=MIST)
tf = box(s, 0.85, 4.18, 11.6, 1.8)
para(tf, "Two operational takeaways", 14.5, TEAL, bold=True, font=F_HDR, space=6)
para(tf, "1.  Evening service is the gap. Measured in-service time collapses after ~19:00 across the observed fleet — "
         "an enforced-headway evening tier would be a visible, cheap service win.", 13, INK, space=6)
para(tf, "2.  The 24-minute turn is what schedulers must plan for. It includes waiting-to-fill, so a scheduled service "
         "can beat it — but any timetable assuming near-zero recovery is fiction.", 13, INK, space=0)

# ════════ 7. HOW IT REFINED THE PLAN (v3.4.5) ═══════════════════════
s = slide()
title_bar(s, "closing the loop", "Measurement corrected the plan: v3.4.5",
          "Five plan routes were confidently matched to observed corridors — enough evidence to re-anchor their cycle times to measured bus physics.")
tf = box(s, 0.55, 1.95, 7.3, 0.4)
para(tf, "Planned one-way time ÷ measured one-way time", 14.5, INK, bold=True, font=F_HDR)
rows = [("Soura – Lal Chowk (FDR-050)", "211 runs", "0.50", "5 → 7"),
        ("Soura – Railway Stn (FDR-262)", "93 runs", "0.55", "4 → 7"),
        ("Srinagar – Pampore (FDR-370)", "79 runs", "0.51", "4 → 5"),
        ("Soura – Qamarwari (FDR-270)", "47 runs", "0.58", "4 → 5"),
        ("Safakadal – Jehangir Chowk (FDR-575)", "40 runs", "0.47", "3 → 3")]
y = 2.45
for name, runs, ratio, fleet in rows:
    card(s, 0.55, y, 7.3, 0.74, fill=MIST)
    tf = box(s, 0.8, y + 0.12, 4.6, 0.5)
    para(tf, name, 12, INK, bold=True, space=0)
    para(tf, runs + " observed", 10, GREY, space=0)
    tf = box(s, 5.5, y + 0.16, 1.1, 0.45)
    para(tf, ratio, 16, RED, bold=True, font=F_HDR, align=PP_ALIGN.CENTER, space=0)
    tf = box(s, 6.7, y + 0.16, 1.0, 0.45)
    para(tf, fleet, 13, TEAL, bold=True, align=PP_ALIGN.CENTER, space=0)
    y += 0.86
c = card(s, 8.2, 1.95, 4.55, 4.85, fill=TEALDK)
tf = box(s, 8.5, 2.2, 3.95, 4.4)
para(tf, "What was hiding it", 15, WHITE, bold=True, font=F_HDR, space=6)
para(tf, "4 of the 5 routes were pinned at the engine's per-km cycle CAP — a safety heuristic that quietly clamped "
         "their cycle times below what buses measurably need.", 12, RGBColor(0xD5, 0xE8, 0xDF), space=8)
para(tf, "The fix: on directly-measured corridors, the measured moving speed replaces the car-model estimate, and the "
         "cap no longer clamps. Only those 5 routes changed.", 12, RGBColor(0xD5, 0xE8, 0xDF), space=10)
para(tf, "Fleet 1,004 → 1,011", 24, WHITE, bold=True, font=F_HDR, space=2)
para(tf, "+7 buses, all on GPS-verified core corridors", 11.5, RGBColor(0x9F, 0xBF, 0xB3), space=0)

# ════════ 8. STOPS ══════════════════════════════════════════════════
s = slide()
title_bar(s, "bonus layer", "An evidence-based stop inventory",
          "Every 40 s – 12 min dwell is a vote for a real stopping place. Clustered, they map where buses actually stop.")
stat(s, 0.55, 1.85, 3.45, "352", "candidate stops", "from 38,431 observed dwell events")
stat(s, 4.2, 1.85, 3.45, "215", "strong stops", "≥3 distinct drivers and ≥10 visits")
stat(s, 0.55, 3.62, 3.45, "30%", "register corroboration", "of official stops in the observed footprint")
stat(s, 4.2, 3.62, 3.45, "135", "enrichment candidates", "observed stops the register doesn't list")
c = card(s, 0.55, 5.45, 7.1, 1.55, fill=MIST)
tf = box(s, 0.8, 5.6, 6.6, 1.3)
para(tf, "Enrichment, not indictment", 13, TEAL, bold=True, font=F_HDR, space=3)
para(tf, "The register is endpoint-oriented, so mid-route stops are “new” by construction. Each candidate ships "
         "with coordinates, visits and driver support — ready for field validation.", 11.5, INK, space=0)
pic(s, "data/stops_vs_register.png", 7.95, 1.55, h=5.45)
caption(s, 7.95, 7.04, 5.0, "Register (blue) vs observed stops — red = candidates.", align=PP_ALIGN.CENTER)

# ════════ 9. LIMITATIONS (dark, honest) ═════════════════════════════
s = slide(TEALDK)
tf = box(s, 0.55, 0.4, 12.2, 1.1)
para(tf, "READ BEFORE CITING", 12, RGBColor(0x7F, 0xC8, 0xAF), bold=True)
para(tf, "What this data cannot tell us", 31, WHITE, bold=True, font=F_HDR, space=2)
lims = [
    ("Partial, self-selected adoption", "~157 active drivers, Srinagar-concentrated. “No app data on a route” does NOT mean “no service”."),
    ("No demand, no ridership, no frequency", "The app sees vehicles, not passengers. Observed run counts reflect who installed the app — never cited as service frequency."),
    ("Corridors cover 41% of clean runs", "The other 59% was decomposed and shown to be noise, loops and out-of-division traffic — but the corridor findings describe the frequently-driven core only."),
    ("Observed times are an upper bound", "They embody informal wait-to-fill practice; a scheduled, headway-enforced service should beat them."),
    ("Rural routes remain unvalidated", "Measured corrections touch 5 Srinagar-core corridors; long rural lifeline cycle times still rest on the model."),
]
y = 1.75
for h, d in lims:
    card(s, 0.55, y, 12.2, 0.92, fill=RGBColor(0x0E, 0x4A, 0x3E))
    tf = box(s, 0.85, y + 0.1, 11.6, 0.74)
    para(tf, h, 13.5, WHITE, bold=True, space=1)
    para(tf, d, 11.5, RGBColor(0xC9, 0xE2, 0xD8), space=0)
    y += 1.06

# ════════ 10. CLOSE ═════════════════════════════════════════════════
s = slide()
title_bar(s, "summary", "A validation layer the plan didn't have before")
left = [("Verified", "7 plan routes confirmed running in reality; 8 geometry divergences queued for reconciliation; 0 informal routes."),
        ("Measured", "First real bus speeds (21 / 12.5 km/h), duty cycles, 24-min turnarounds and a 10,600-cell congestion map for Srinagar."),
        ("Corrected", "5 GPS-verified corridors re-anchored to measured physics → plan v3.4.5 (1,011 buses), with the audit trail published.")]
y = 1.7
for h, d in left:
    card(s, 0.55, y, 7.5, 1.34, fill=MIST)
    tf = box(s, 0.85, y + 0.16, 6.9, 1.05)
    para(tf, h, 15, TEAL, bold=True, font=F_HDR, space=2)
    para(tf, d, 12.5, INK, space=0)
    y += 1.5
c = card(s, 8.4, 1.7, 4.35, 3.15, fill=TEALDK)
tf = box(s, 8.7, 1.95, 3.75, 2.75)
para(tf, "To make this decisive", 15, WHITE, bold=True, font=F_HDR, space=8)
para(tf, "1.  Broader driver adoption — every new driver hardens every layer.", 12.5, RGBColor(0xD5, 0xE8, 0xDF), space=6)
para(tf, "2.  The RTO's surveyed stop register — to confirm the 135 candidates.", 12.5, RGBColor(0xD5, 0xE8, 0xDF), space=6)
para(tf, "3.  AFC / ticketing data — the missing demand half no GPS can supply.", 12.5, RGBColor(0xD5, 0xE8, 0xDF), space=0)
tf = box(s, 0.55, 6.55, 12.2, 0.6)
para(tf, "Repo: github.com/Princu-Babu/bus-sathi-trace-intelligence  ·  full method, audit log and reality-check tables committed",
     11.5, GREY, italic=True)

prs.save("Bus_Sathi_Trace_Intelligence_Briefing.pptx")
print(f"Wrote Bus_Sathi_Trace_Intelligence_Briefing.pptx ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
